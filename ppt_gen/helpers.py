"""Layout helpers for python-pptx slides."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ppt_gen import theme as T


def set_fonts(run, latin: str = T.LATIN_FONT, ea: str = T.EA_FONT) -> None:
    """Set both latin and east-asian typefaces so CJK renders correctly."""
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea_el)
    ea_el.set("typeface", ea)


def style_run(run, size: int, color: RGBColor = T.INK, bold: bool = False,
              mono: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_fonts(run, latin=T.MONO_FONT if mono else T.LATIN_FONT,
              ea=T.EA_FONT if not mono else T.MONO_FONT)


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text: str, size: int = 16, color: RGBColor = T.INK,
         bold: bool = False, level: int = 0, space_after: int = 6,
         bullet: bool = False, mono: bool = False, align=None, first: bool = False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    style_run(run, size, color, bold, mono)
    return p


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = T.SKY
    bar.line.fill.background()

    tf = textbox(slide, Inches(0.8), Inches(0.3), Inches(11.9), Inches(1.0))
    para(tf, text, size=28, color=T.NAVY, bold=True, first=True, space_after=2)
    if subtitle:
        para(tf, subtitle, size=13, color=T.MUTED)


def card(slide, x, y, w, h, fill: RGBColor = T.CARD_BG, line: RGBColor | None = None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def picture(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def caption(slide, text: str, x, y, w, size: int = 12,
            color: RGBColor = T.MUTED, align=PP_ALIGN.CENTER):
    tf = textbox(slide, x, y, w, Inches(0.35))
    para(tf, text, size=size, color=color, first=True, align=align)


def footer(slide, idx: int, total: int) -> None:
    tf = textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.35))
    para(tf, "KFB Cloud Coverage", size=10, color=T.MUTED, first=True)
    tf2 = textbox(slide, Inches(12.2), Inches(7.05), Inches(0.7), Inches(0.35))
    para(tf2, f"{idx} / {total}", size=10, color=T.MUTED, first=True,
         align=PP_ALIGN.RIGHT)


def chevron(slide, x, y, w, h, text: str, sub: str | None = None,
            fill: RGBColor = T.SKY) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(60000)
    para(tf, text, size=13, color=T.WHITE, bold=True, first=True,
         align=PP_ALIGN.CENTER, space_after=0)
    if sub:
        para(tf, sub, size=10, color=T.SKY_LIGHT, align=PP_ALIGN.CENTER,
             space_after=0)
