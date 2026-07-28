"""Slide builders — one function per slide."""
from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ppt_gen import theme as T
from ppt_gen.helpers import (
    add_title, caption, card, chevron, footer, para, picture, textbox,
)

TOTAL = 11


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s01_cover(prs):
    s = _blank(prs)
    picture(s, T.ASSETS / "cover.jpg", 0, 0, w=T.SLIDE_W, h=T.SLIDE_H)
    tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.4))
    para(tf, "KFB Cloud Coverage", size=48, color=T.WHITE, bold=True,
         first=True, space_after=6)
    para(tf, "Estimating cloud base from a hill weather-camera photo",
         size=26, color=T.SKY_LIGHT, space_after=18)
    para(tf, "One Kadoorie Farm photo (camera at ~150 m) → day/night → inside cloud? → cloud-base message",
         size=15, color=T.WHITE)
    tf2 = textbox(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5))
    para(tf2, "Project overview · July 2026", size=13, color=T.SKY_LIGHT, first=True)


def s02_problem(prs):
    s = _blank(prs)
    add_title(s, "Problem & Motivation",
              "Drone pilots need cloud-base height, but Hong Kong has few direct observations")
    rows = [
        ("Regulatory need",
         "CAD drone rules require pilots to know the approximate cloud-base height."),
        ("Observation gap",
         "Hong Kong has almost no public, real-time direct measurements of cloud base."),
        ("Existing asset",
         "HKO operates hill weather cameras at known altitudes — e.g. Kadoorie Farm (KFB) at ~150 m."),
    ]
    y = Inches(1.55)
    for title, body in rows:
        card(s, Inches(0.7), y, Inches(7.0), Inches(1.35))
        tf = textbox(s, Inches(1.0), y + Inches(0.16), Inches(6.5), Inches(1.1))
        para(tf, title, size=16, color=T.SKY, bold=True, first=True, space_after=3)
        para(tf, body, size=14, color=T.INK)
        y += Inches(1.6)

    card(s, Inches(8.0), Inches(1.55), Inches(4.7), Inches(4.55), fill=T.CARD_SKY)
    tf = textbox(s, Inches(8.3), Inches(1.8), Inches(4.1), Inches(4.2))
    para(tf, "Key inference", size=17, color=T.NAVY, bold=True, first=True, space_after=8)
    para(tf, "Camera is inside cloud (whiteout, distant view lost)",
         size=14, bold=True, space_after=2)
    para(tf, "→ cloud base ≤ 150 m (camera altitude)", size=14, color=T.SKY,
         space_after=12)
    para(tf, "Camera is not inside cloud (distant view clear)",
         size=14, bold=True, space_after=2)
    para(tf, "→ cloud base > 150 m", size=14, color=T.GREEN, space_after=12)
    para(tf, "Turn an existing webcam into a cloud-base sensor — no new hardware.",
         size=13, color=T.MUTED)


def s03_idea(prs):
    s = _blank(prs)
    add_title(s, "Core Idea",
              "No radar, no new hardware — just the visual cues in the photo")
    specs = [
        ("Daytime cues", "imgKFB_160106_1115.jpg", "imgKFB_160105_1150.jpg",
         "Clear: valley, towns and distant hills all visible",
         "In cloud: uniform whiteout, distant texture gone"),
        ("Nighttime cues", "imgKFB_160101_0100.jpg", "imgKFB_160103_0320.jpg",
         "Clear: valley lights clearly visible",
         "In cloud: lights swallowed by fog, uniformly dark frame"),
    ]
    x = Inches(0.7)
    for title, clear_img, fog_img, clear_cap, fog_cap in specs:
        card(s, x, Inches(1.5), Inches(6.0), Inches(5.15))
        tf = textbox(s, x + Inches(0.3), Inches(1.68), Inches(5.4), Inches(0.5))
        para(tf, title, size=17, color=T.NAVY, bold=True, first=True)
        pic_w = Inches(5.4)
        pic_h = Inches(1.72)
        picture(s, T.IMAGES / clear_img, x + Inches(0.3), Inches(2.25),
                w=pic_w, h=pic_h)
        caption(s, clear_cap, x + Inches(0.3), Inches(4.0), pic_w, size=12,
                color=T.GREEN, align=PP_ALIGN.LEFT)
        picture(s, T.IMAGES / fog_img, x + Inches(0.3), Inches(4.42),
                w=pic_w, h=pic_h)
        caption(s, fog_cap, x + Inches(0.3), Inches(6.17), pic_w, size=12,
                color=T.SKY, align=PP_ALIGN.LEFT)
        x += Inches(6.35)


def s04_pipeline(prs):
    s = _blank(prs)
    add_title(s, "System Pipeline",
              "Classical image features + a lightweight model, fully automated end to end")
    steps = [
        ("Image input", "3,497 KFB photos"),
        ("Feature extraction", "brightness / edges / saturation…"),
        ("Heuristic auto-label", "labels.csv, human-editable"),
        ("Train classifier", "logistic regression + scaler"),
        ("Predict", "fixed cloud-base message"),
    ]
    x = Inches(0.55)
    w = Inches(2.55)
    for text, sub in steps:
        chevron(s, x, Inches(1.7), w, Inches(1.05), text, sub)
        x += Inches(2.47)

    card(s, Inches(0.7), Inches(3.2), Inches(12.0), Inches(3.5))
    tf = textbox(s, Inches(1.0), Inches(3.4), Inches(11.4), Inches(3.2))
    para(tf, "Module layout (each file does one job)", size=16, color=T.NAVY,
         bold=True, first=True, space_after=8)
    files = [
        ("src/features.py",
         "load image, crop the timestamp banner, extract 10 features, decide day/night"),
        ("src/heuristic.py",
         "visual rules produce seed labels (inside_cloud / not_inside)"),
        ("src/auto_label.py",
         "scan every photo in image/, write data/labels.csv"),
        ("src/train.py",
         "80/20 split by date, train logistic regression, export a tiny JSON model"),
        ("src/predict.py",
         "single/batch inference, print the fixed message"),
        ("src/web.py + docs/",
         "Flask server variant and the GitHub Pages in-browser variant"),
    ]
    for name, desc in files:
        p = tf.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = name
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.name = T.MONO_FONT
        r1.font.color.rgb = T.SKY
        r2 = p.add_run()
        r2.text = "   " + desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = T.INK
        from ppt_gen.helpers import set_fonts
        set_fonts(r2)


def s05_features(prs):
    s = _blank(prs)
    add_title(s, "Feature Engineering",
              "Crop the timestamp banner (top 8%), resize to 640 px long side — identical to the browser Canvas")
    groups = [
        ("Overall appearance", [
            "Brightness mean / std (whiteout → high mean, low variance)",
            "Saturation mean (inside cloud looks grey, low saturation)",
            "Edge density (Laplacian; little texture inside cloud)"]),
        ("Spatial structure", [
            "Upper vs lower half contrast (is the far scene washed out?)",
            "Far-field gradient far_grad (excludes bottom 28% trees)",
            "Far-field washout ratio far_wash and std far_std"]),
        ("Night-specific", [
            "Bright-spot ratio bright_spot_ratio (valley lights)",
            "No lights + uniform frame → lean inside_cloud"]),
        ("Day/night decision", [
            "Image brightness only: mean ≥ 80 means day",
            "No dependence on filename time → works for any photo source"]),
    ]
    coords = [(0.7, 1.7), (6.4, 1.7), (0.7, 4.25), (6.4, 4.25)]
    for (gx, gy), (title, items) in zip(coords, groups):
        x, y = Inches(gx), Inches(gy)
        card(s, x, y, Inches(5.45), Inches(2.3),
             fill=T.CARD_SKY if title == "Night-specific" else T.CARD_BG)
        tf = textbox(s, x + Inches(0.28), y + Inches(0.18), Inches(4.9), Inches(2.0))
        para(tf, title, size=15, color=T.NAVY, bold=True, first=True, space_after=5)
        for it in items:
            para(tf, it, size=12.5, bullet=True, space_after=3)


def s06_data(prs):
    s = _blank(prs)
    add_title(s, "Data & Auto-Labeling",
              "Heuristic rules label everything first; humans only spot-check — no hand-labeling 3,497 photos")
    picture(s, T.ASSETS / "distribution.png", Inches(0.55), Inches(1.75),
            w=Inches(6.6))
    card(s, Inches(7.5), Inches(1.7), Inches(5.25), Inches(4.9))
    tf = textbox(s, Inches(7.8), Inches(1.92), Inches(4.65), Inches(4.5))
    para(tf, "Labeling workflow", size=16, color=T.NAVY, bold=True, first=True,
         space_after=6)
    para(tf, "python -m src.auto_label scans all photos and applies the visual rules as seed labels",
         size=13, bullet=True, space_after=4)
    para(tf, "Output goes to data/labels.csv (with all feature values; edit the label column to correct)",
         size=13, bullet=True, space_after=4)
    para(tf, "src/sample_review.py prints stratified samples (day/night × in/out of cloud) for an eye check",
         size=13, bullet=True, space_after=4)
    para(tf, "Retrain after corrections — labeling cost stays tiny",
         size=13, bullet=True, space_after=10)
    para(tf, "Example heuristic rule (daytime)", size=15, color=T.NAVY, bold=True,
         space_after=4)
    para(tf, "far_wash >= 0.60 and far_grad <= 5.0", size=12, mono=True,
         space_after=2)
    para(tf, "→ far field mostly washed out with almost no texture → inside cloud",
         size=12.5, color=T.SKY)


def s07_model(prs):
    s = _blank(prs)
    add_title(s, "Model & Evaluation",
              "Logistic regression (StandardScaler + class_weight=balanced), split by date to avoid same-day leakage")
    picture(s, T.ASSETS / "confusion.png", Inches(0.6), Inches(1.75), h=Inches(4.1))
    picture(s, T.ASSETS / "features.png", Inches(5.6), Inches(1.75), h=Inches(4.1))
    card(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.85), fill=T.CARD_GREEN)
    tf = textbox(s, Inches(1.0), Inches(6.2), Inches(11.4), Inches(0.6))
    para(tf, "Test accuracy 92.4% (458 images from held-out dates). Weights match intuition: "
             "high brightness variance, more edges and higher saturation all push toward not_inside.",
         size=14, color=T.INK, first=True)


def s08_demo(prs):
    s = _blank(prs)
    add_title(s, "Live Examples",
              "Real outputs for four scenarios (python -m src.predict)")
    cases = [
        ("imgKFB_160106_1115.jpg", "Day · clear",
         "I am not inside cloud. Cloud base above 150 m.", T.GREEN),
        ("imgKFB_160105_1150.jpg", "Day · in cloud",
         "Camera is inside cloud (day). Cloud base at or below ~150 m.", T.SKY),
        ("imgKFB_160101_0100.jpg", "Night · clear",
         "I am not inside cloud. Cloud base above 150 m.", T.GREEN),
        ("imgKFB_160103_0320.jpg", "Night · in cloud",
         "Camera is inside cloud (night). Cloud base at or below ~150 m.", T.SKY),
    ]
    coords = [(0.7, 1.55), (7.0, 1.55), (0.7, 4.3), (7.0, 4.3)]
    for (cx, cy), (img, tag, msg, color) in zip(coords, cases):
        x, y = Inches(cx), Inches(cy)
        picture(s, T.IMAGES / img, x, y, w=Inches(3.4), h=Inches(1.91))
        tf = textbox(s, x + Inches(3.55), y + Inches(0.1), Inches(2.55), Inches(2.2))
        para(tf, tag, size=15, color=T.NAVY, bold=True, first=True, space_after=5)
        para(tf, msg, size=11.5, color=color, mono=True)


def s09_deploy(prs):
    s = _blank(prs)
    add_title(s, "Deployment Strategy",
              "One rule set, three runtimes, identical results")
    cols = [
        ("GitHub Pages (primary)", T.CARD_SKY, [
            "Static site in docs/; app.js computes in the visitor's browser",
            "Photos never leave the user's device → privacy-friendly",
            "Zero server cost, no OOM — anyone with the link can use it",
        ]),
        ("Render / Flask (fallback)", T.CARD_BG, [
            "src/web.py: Flask + waitress, single thread to save memory",
            "BLAS threads and upload size capped (8 MB) for the free tier",
            "Dockerfile included — portable to any container platform",
        ]),
        ("CLI (development)", T.CARD_BG, [
            "python -m src.predict photo.jpg prints the message directly",
            "auto_label / train / sample_review are all one-liners",
            "Model exports to a tiny JSON (weights + intercept) readable in JS",
        ]),
    ]
    x = Inches(0.7)
    for title, fill, items in cols:
        card(s, x, Inches(1.6), Inches(3.95), Inches(4.9), fill=fill)
        tf = textbox(s, x + Inches(0.28), Inches(1.85), Inches(3.4), Inches(4.4))
        para(tf, title, size=16, color=T.NAVY, bold=True, first=True, space_after=8)
        for it in items:
            para(tf, it, size=13, bullet=True, space_after=8)
        x += Inches(4.15)
    tf = textbox(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4))
    para(tf, "The Python and JavaScript feature code is deliberately aligned line by line "
             "(same resize, crop and formulas), so all three runtimes agree.",
         size=13, color=T.MUTED, first=True)


def s10_engineering(prs):
    s = _blank(prs)
    add_title(s, "Engineering Practices", "A small project run with a proper process")
    items = [
        ("Test-driven development",
         "Failing tests first, then implementation: message format, day/night, feature shape and "
         "heuristic rules each have unit tests, plus an end-to-end smoke test on real images."),
        ("No data leakage",
         "Train/test split is by date, so photos from the same day never appear on both sides — "
         "frames five minutes apart are nearly identical."),
        ("Spec and plan first",
         "docs/superpowers/ holds the design spec (goals, non-goals, success criteria) and a "
         "step-by-step implementation plan, reviewed before coding."),
        ("Deliberate non-goals",
         "v1 skips CNNs, humidity fusion, multi-site cameras and continuous cloud-base regression — "
         "make the simplest version useful first."),
    ]
    y = Inches(1.6)
    for title, body in items:
        card(s, Inches(0.7), y, Inches(12.0), Inches(1.15))
        tf = textbox(s, Inches(1.0), y + Inches(0.12), Inches(11.4), Inches(0.95))
        para(tf, title, size=15, color=T.SKY, bold=True, first=True, space_after=2)
        para(tf, body, size=13, color=T.INK)
        y += Inches(1.32)


def s11_future(prs):
    s = _blank(prs)
    add_title(s, "Summary & Future Work")
    card(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(4.9), fill=T.CARD_GREEN)
    tf = textbox(s, Inches(1.0), Inches(1.85), Inches(5.4), Inches(4.4))
    para(tf, "Done", size=17, color=T.NAVY, bold=True, first=True, space_after=8)
    for it in [
        "Fully automated pipeline: 3,497 photos → labels → training → inference",
        "Works day and night, emitting two fixed cloud-base messages",
        "92.4% test accuracy, validated on held-out dates",
        "Live on GitHub Pages — anyone can use it, photos stay on-device",
        "Unit tests + smoke test all green",
    ]:
        para(tf, it, size=13.5, bullet=True, space_after=7)

    card(s, Inches(7.0), Inches(1.6), Inches(5.75), Inches(4.9), fill=T.CARD_SKY)
    tf = textbox(s, Inches(7.3), Inches(1.85), Inches(5.15), Inches(4.4))
    para(tf, "Future work (v2+)", size=17, color=T.NAVY, bold=True, first=True,
         space_after=8)
    for it in [
        "Fuse relative humidity (RH) and other weather data for reliability",
        "Add more sites: Tai Mo Shan, Victoria Peak — cameras at other altitudes",
        "Combine altitudes to upgrade from a binary call to a cloud-base height range",
        "Consider a lightweight CNN only if classical features plateau on corrected labels",
        "Collect data across seasons to handle rain and haze confounders",
    ]:
        para(tf, it, size=13.5, bullet=True, space_after=7)

    tf = textbox(s, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.5))
    para(tf, "One photo, one answer: is the cloud base above or below 150 metres?",
         size=15, color=T.NAVY, bold=True, first=True, align=PP_ALIGN.CENTER)


BUILDERS = [s01_cover, s02_problem, s03_idea, s04_pipeline, s05_features,
            s06_data, s07_model, s08_demo, s09_deploy, s10_engineering,
            s11_future]


def build_slides(prs) -> None:
    for i, fn in enumerate(BUILDERS, start=1):
        fn(prs)
        if i > 1:
            footer(prs.slides[-1], i, TOTAL)
