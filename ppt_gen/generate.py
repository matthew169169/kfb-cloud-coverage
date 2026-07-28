"""Entry point: build assets then write the deck.

Usage: python -m ppt_gen.generate
"""
from pptx import Presentation

from ppt_gen import theme as T
from ppt_gen.charts import build_all
from ppt_gen.slides import build_slides

OUT = T.ROOT / "KFB_Cloud_Coverage.pptx"


def main() -> None:
    build_all()
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    build_slides(prs)
    prs.save(OUT)
    print(f"saved {OUT} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
